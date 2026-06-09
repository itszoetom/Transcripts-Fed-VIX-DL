"""Rewrite the three model boxes on slide 4 of Presentation_D3.pptx.

Goal: more technical wording that keeps the proper terms (encoder, embedding,
masked language modeling, additive/Bahdanau attention, softmax, linear
regression head, gradient descent, MSE) but glosses each one in plain words.

Each box is rebuilt from scratch (title kept, body paragraphs replaced) so it
is robust to whatever manual edits the box currently holds. Shrink-to-fit is
enabled so PowerPoint keeps the denser text inside the existing box outlines.
Only slide 4 is touched; every other slide is left exactly as-is.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

PATH = Path("presentation/Presentation_D3.pptx")
BODY_PT = 11

BOXES = {
    "Rounded Rectangle 6": (  # FinBERT
        "FinBERT (frozen encoder)",
        [
            "Encoder: a BERT transformer that maps each sentence to a 768-d embedding (a vector capturing its meaning).",
            "Pre-trained on 4.9 billion words of finance text by masked language modeling (learning by predicting hidden words).",
            "Frozen: its 110M weights never update. 198 documents is far too few to fine-tune that many parameters.",
        ],
    ),
    "Rounded Rectangle 11": (  # attention
        "Sentence attention (trained)",
        [
            "Additive (Bahdanau) attention: a small trained network gives each of the 80 sentence embeddings a relevance score.",
            "A softmax converts the 80 scores into weights that sum to 1 (a probability distribution over the sentences).",
            "The document embedding is their weighted sum, so high-weight sentences dominate. The ~99k weights are trained and interpretable.",
        ],
    ),
    "Rounded Rectangle 18": (  # linear head
        "Linear regression head (trained)",
        [
            "Maps the 768-d document embedding to a single scalar (one number): the predicted 10-day VIX change.",
            "y = w*d + b   (a weighted sum of the 768 features plus a bias term).",
            "Trained jointly with the attention by gradient descent on MSE loss (minimizes average squared error).",
            "Just one linear layer, kept minimal to avoid overfitting on small data.",
        ],
    ),
}


def rebuild(tf, title, items):
    body = tf._txBody
    paras = body.findall(qn("a:p"))

    # Keep the first paragraph as the title; overwrite its text but keep its
    # run formatting (16pt bold).
    title_para = tf.paragraphs[0]
    if title_para.runs:
        title_para.runs[0].text = title
        for r in title_para.runs[1:]:
            r.text = ""
    else:
        title_para.add_run().text = title

    # Drop every paragraph after the title.
    for p in paras[1:]:
        body.remove(p)

    # Blank spacer line under the title (matches the original look).
    tf.add_paragraph()

    for text in items:
        para = tf.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = text
        run.font.size = Pt(BODY_PT)

    # Let PowerPoint shrink the font if the denser text overflows the box.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = True


def main():
    prs = Presentation(str(PATH))
    slide = prs.slides[3]
    by_name = {s.name: s for s in slide.shapes}
    for name, (title, items) in BOXES.items():
        rebuild(by_name[name].text_frame, title, items)
    prs.save(str(PATH))
    print(f"updated model boxes in {PATH}")


if __name__ == "__main__":
    main()
