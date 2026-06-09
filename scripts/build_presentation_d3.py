"""Patch Presentation_D2.pptx into Presentation_D3.pptx.

D3 changes (content only, all formatting preserved by editing runs in place):
  - Slide 4 (model): simpler, freshman-level descriptions in the FinBERT,
    attention, and linear-head boxes. FinBERT numbers corrected (trained on
    4.9 billion words; ~110 million weights). "233" updated to the real
    training-set size of 198.
  - Slide 5 (results): drop the validation result (not truly out of sample);
    honest per-regime framing (nothing significant; skill rises with recency,
    the opposite of the decay we expected; beats baseline in 2 of 3 regimes).
  - Slide 6 (conclusion): remove the 3-day-vs-10-day horizon story and the
    "became significant" claim; refocus on the results we actually showed and
    on how the model worked or did not.

The slide-3 split diagram already reads TRAIN 198 / VAL 35, so it is untouched.
"""

from pathlib import Path
from pptx import Presentation

SRC = Path("presentation/Presentation_D2.pptx")
OUT = Path("presentation/Presentation_D3.pptx")


def shape_by_name(slide, name):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    raise KeyError(f"shape {name!r} not found on slide")


def set_para(tf, idx, text):
    """Replace the visible text of paragraph `idx`, keeping its run formatting.

    Writes into the first run and blanks any trailing runs so the original
    font size / colour / bullet level are preserved.
    """
    para = tf.paragraphs[idx]
    runs = para.runs
    if not runs:
        raise ValueError(f"paragraph {idx} has no runs to format from: {para.text!r}")
    runs[0].text = text
    for r in runs[1:]:
        r.text = ""


def main():
    prs = Presentation(str(SRC))
    s_model = prs.slides[3]
    s_results = prs.slides[4]
    s_concl = prs.slides[5]

    # ---- Slide 4: model boxes -------------------------------------------
    finbert = shape_by_name(s_model, "Rounded Rectangle 6").text_frame
    set_para(finbert, 2, "Trained on 4.9 billion words of finance text, so it reads Fed language better than plain BERT.")
    set_para(finbert, 3, "It has about 110 million weights, and we froze all of them, so it just turns each sentence into numbers.")
    set_para(finbert, 4, "We froze it because 198 documents is far too few to safely retrain 110 million weights.")

    attn = shape_by_name(s_model, "Rounded Rectangle 11").text_frame
    set_para(attn, 0, "Attention layer (trained)")
    set_para(attn, 2, "For each of the 80 sentences, it learns an importance score for how much that sentence matters.")
    set_para(attn, 3, "It then merges the sentences into one summary, letting the important ones count more.")
    set_para(attn, 4, "We can read those scores to see which sentences the model focused on.")

    head = shape_by_name(s_model, "Rounded Rectangle 18").text_frame
    set_para(head, 2, "Takes the document summary and turns it into one number: the predicted VIX change.")
    # P3 keeps the equation "y = w*d + b"
    set_para(head, 5, "It is just a weighted sum plus an offset, the simplest possible predictor.")
    set_para(head, 6, "We kept it this simple on purpose so it would not overfit our small dataset.")

    # ---- Slide 5: results -----------------------------------------------
    res = shape_by_name(s_results, "TextBox 10").text_frame
    # P0-P2 (metric) and P4 ("Results") header stay; rewrite the four result bullets,
    # reusing the old validation-line slot so no paragraph is added or removed.
    set_para(res, 5, "No regime is statistically significant (all p > 0.05): no reliable out-of-sample prediction.")
    set_para(res, 6, "Skill changes across regimes: weakest in regime 1 (r = +0.07), strongest in regime 3 (r = +0.33).")
    set_para(res, 7, "Opposite of what we expected: it did best on the most recent era, furthest from training.")
    set_para(res, 8, "Beats the TF-IDF baseline in regimes 2 and 3, but loses to it in regime 1.")

    caveat = shape_by_name(s_results, "TextBox 12").text_frame
    set_para(caveat, 0, "Caveat: regime 3 has only n = 13 documents (regimes 1 and 2 have 40 each).")
    set_para(caveat, 1, "That r = +0.33 is not significant (p = 0.27) and has a wide confidence interval.")
    set_para(caveat, 2, "So the strong recent result is most likely small-sample noise, not a real trend.")

    # ---- Slide 6: conclusion --------------------------------------------
    concl = shape_by_name(s_concl, "Content Placeholder 2").text_frame
    set_para(concl, 0, "What we found")
    set_para(concl, 1, "The model did not reliably predict VIX change: no test regime was statistically significant.")
    set_para(concl, 2, "Its skill changed across regimes, doing best on the most recent era, the opposite of what we expected.")
    set_para(concl, 3, "Frozen FinBERT plus a small attention layer was a sensible design for so little data, and the attention weights showed which sentences mattered.")
    set_para(concl, 4, "The simple TF-IDF baseline was competitive, so the deep model is not clearly better yet.")
    set_para(concl, 5, "What we would do differently")
    set_para(concl, 6, "FOMC statements (released same day) instead of minutes (released about 3 weeks later).")
    set_para(concl, 7, "An intraday VIX target to capture the immediate market reaction.")
    set_para(concl, 8, "More data, especially to firm up the recent regimes.")

    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
